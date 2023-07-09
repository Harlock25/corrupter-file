import tkinter as tk
from tkinter import *
from tkinter import messagebox
from tkinter import filedialog
import os
import re 
import shutil
from pptx import Presentation
import codecs
import aspose.words as aw
from openpyxl import Workbook
from openpyxl.worksheet.table import Table, TableStyleInfo

class StartWindow(tk.Toplevel):
    def __init__(self, parent):
        super().__init__(parent)
        self.title("Start Window")
        self.attributes('-fullscreen', True)
        self.fullScreenState = False
        self.bind("<F11>", self.toggleFullScreen)
        self.bind("<Escape>", self.quitFullScreen)
        self.config(background="red")

        self.label_title = tk.Label(self, text="FILE CORRUPTOR", font=("Arial", 45, "bold"),bg="red", fg="white")
        self.label_title.grid(row=0, column=0, columnspan=5, pady=15, sticky="ew")

        self.grid_columnconfigure(0, weight=5)
        self.label_title.grid(sticky="ew")

        self.label_file = tk.Label(self, text="Filename (format: name.ext)", font=("Arial", 16), bg="red", fg="white")
        self.label_file.grid(row=3, column=0, padx=10, pady=150, sticky="e")

        self.input_area_filename = tk.Entry(self, font=("Arial", 16))
        self.input_area_filename.grid(row=3, column=2, padx=10, pady=150, sticky="w")

        self.label_point = tk.Label(self, text=".", font=("Arial", 16), bg="red", fg="white")
        self.label_point.grid(row=3, column=3, padx=10, pady=150, sticky="e")

        self.input_extension = tk.Entry(self, font=("Arial", 16))
        self.input_extension.grid(row=3, column=4, padx=10, pady=150, sticky="w")

        self.grid_columnconfigure(0, weight=3)
        self.grid_columnconfigure(4, weight=3)

        open_button = tk.Button(self, text="CREATE FILE", height=3, width=15, command=self.create_file, font=("Arial", 18), bg="blue", fg="white")
        open_button.grid(row=4, column=2, pady=10)

        open_button = tk.Button(self, text="OPEN", height=3, width=15, command=self.open_file, font=("Arial", 18), bg="green", fg="white")
        open_button.grid(row=5, column=2, pady=10)

    def toggleFullScreen(self, event):
        self.fullScreenState = not self.fullScreenState
        self.attributes("-fullscreen", self.fullScreenState)

    def quitFullScreen(self, event):
        self.fullScreenState = False
        self.attributes("-fullscreen", self.fullScreenState)

    def validate_name(self, name):
        return bool(re.match(r'^[a-zA-Z0-9_.]+$', name))

    def validate_extension(self, extension):
        return bool(re.match(r'^[a-zA-Z0-9_.]+$', extension))

    def create_file(self):
        name = self.input_area_filename.get()
        extension = self.input_extension.get()

        if not self.validate_name(name) or not self.validate_extension(extension):
            messagebox.showerror("Invalid format", "Please enter a valid filename.")
            return

        filename = name + "." + extension
        file_path = os.path.join("./files", filename)

        try:
            if extension == "pptx":
                prs = Presentation()
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                title.text = "Hey, This is a Slide! How exciting!"
                subtitle.text = "Really?"
                prs.save(file_path)
            elif extension == "docx":
                doc = aw.Document()
                builder = aw.DocumentBuilder(doc)
                builder.writeln("Sample text")
                doc.save(file_path)
            elif extension == "xls":
                wb = Workbook()
                ws = wb.active
                data = [
                    ['Apples', 10000, 5000, 8000, 6000],
                    ['Pears',   2000, 3000, 4000, 5000],
                    ['Bananas', 6000, 6000, 6500, 6000],
                    ['Oranges',  500,  300,  200,  700],
                ]
                ws.append(["Fruit", "2011", "2012", "2013", "2014"])
                for row in data:
                    ws.append(row)
                tab = Table(displayName="Table1", ref="A1:E5")
                style = TableStyleInfo(name="TableStyleMedium9", showFirstColumn=False,
                                       showLastColumn=False, showRowStripes=True, showColumnStripes=True)
                tab.tableStyleInfo = style
                ws.add_table(tab)
                wb.save(file_path)
            else:
                messagebox.showerror("Format not valid", "The file format is not supported.")

            messagebox.showinfo("Success", "File created successfully.")
        except Exception as e:
            messagebox.showerror("Error", str(e))

    def open_file(self):
        filetypes = (
            ('text files', '*.txt'),
            ('All files', '*.*')
        )

        filenames = filedialog.askopenfilenames(
            title='Open files',
            initialdir='/',
            filetypes=filetypes)

        for filepath in filenames:
            file_name = os.path.splitext(os.path.basename(filepath))[0]
            file_extension = os.path.splitext(filepath)[1]
            print("Selected file:", file_name)
            print("File extension:", file_extension)
            self.corrupt_file_open(filepath)

    def corrupt_file_open(self, filepath):
        filename_new = os.path.basename(filepath)
        file_path_new = os.path.join("./files", filename_new)
        os.rename(filepath, file_path_new)

        with codecs.open(file_path_new, 'r', encoding='ansi') as file:
            content = file.read()

        with codecs.open(file_path_new, 'w', encoding='utf-8') as file:
            file.write(content)

        messagebox.showinfo("Success", "File renamed and re-encoded successfully.")
